from pathlib import Path
import argparse

def generate_sample_ppt(output_path: str = "static/samples/sample_lecture.pptx") -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        raise SystemExit(f"python-pptx is required: pip install python-pptx ({e})")

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Introduction to Machine Learning"
    s1.placeholders[1].text = "An overview of how machines learn patterns from data."

    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Key Concepts"
    tf = s2.shapes.placeholders[1].text_frame
    tf.text = "Supervised Learning"
    p = tf.add_paragraph(); p.text = "Unsupervised Learning"; p.level = 1
    p = tf.add_paragraph(); p.text = "Reinforcement Learning"; p.level = 1

    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "Algorithms"
    rows, cols = 3, 3
    table = s3.shapes.add_table(rows, cols, Inches(1.0), Inches(2.0), Inches(8.0), Inches(2.0)).table
    headers = ["Type", "Example", "Use Case"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    table.cell(1, 0).text = "Supervised"; table.cell(1, 1).text = "Linear Regression"; table.cell(1, 2).text = "Prediction"
    table.cell(2, 0).text = "Unsupervised"; table.cell(2, 1).text = "K-Means"; table.cell(2, 2).text = "Clustering"

    prs.save(output_path)
    print(f"Sample PPT generated at: {output_path}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="Generate a sample PPT for testing")
    ap.add_argument("--out", dest="out", default="static/samples/sample_lecture.pptx")
    args = ap.parse_args()
    generate_sample_ppt(args.out)
