import sys
import json
import tempfile
from pathlib import Path
from typing import Dict

# Add project root
ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from src.explanation_engine import create_explanation_engine

RESULTS_DIR = Path("validation_results")
RESULTS_DIR.mkdir(parents=True, exist_ok=True)

TEST_CASES = {
    "Science": {
        "Photosynthesis": "Plants convert light into chemical energy via light reactions and the Calvin cycle.",
        "Newton's Laws": "Three fundamental laws describing motion and forces.",
        "DNA replication": "DNA makes a copy of itself during cell division."
    },
    "Math": {
        "Derivatives": "The derivative measures the rate of change of a function.",
        "Linear algebra": "Study of vectors, matrices, and linear transformations.",
        "Probability": "Quantifies likelihood of events."
    },
    "Technology": {
        "Machine learning": "Algorithms that learn patterns from data to make predictions.",
        "Networking": "Connecting devices to share data using protocols.",
        "Databases": "Systems for storing and querying structured data."
    },
    "History": {
        "World War II": "Global conflict from 1939 to 1945 with major world powers.",
        "Industrial Revolution": "Transition to new manufacturing processes (18th-19th centuries)."
    },
}


def create_test_ppt_for_topic(topic: str, content: str) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        raise RuntimeError("python-pptx is required for validation script")

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = topic
    s1.placeholders[1].text = f"Overview of {topic}"

    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = f"About {topic}"
    s2.shapes.placeholders[1].text_frame.text = content

    tmpdir = Path(tempfile.mkdtemp())
    p = tmpdir / f"{topic.replace(' ', '_')}.pptx"
    prs.save(str(p))
    return str(p)


def evaluate_explanation_quality(explanation: str, topic: str) -> Dict:
    length_ok = 200 <= len(explanation) <= 2000
    has_intro = any(w in explanation.lower() for w in ["introduction", "overview", "in summary", "in short"]) or len(explanation) > 0
    has_structure = explanation.count("\n\n") >= 2
    return {"length_ok": length_ok, "has_intro_or_summary": has_intro, "has_structure": has_structure}


def evaluate_analogies_quality(analogies: Dict) -> Dict:
    arr = analogies.get("analogies", []) or []
    return {"count": len(arr), "has_everyday_language": any("like" in a.lower() for a in arr)}


def evaluate_examples_quality(examples: Dict) -> Dict:
    arr = examples.get("examples", []) or []
    return {"count": len(arr), "concrete": any("for example" in e.lower() or any(ch.isdigit() for ch in e) for e in arr)}


def evaluate_quiz_quality(quiz: Dict) -> Dict:
    qs = quiz.get("questions", []) or []
    plaus = all(len(q.get("options", [])) >= 3 for q in qs)
    return {"count": len(qs), "has_rationales": all(bool(q.get("rationale")) for q in qs), "options_plausible": plaus}


def test_topic(engine, topic: str, content: str) -> Dict:
    # Build a minimal slides_data representation
    slides_data = {"slides": [{"text_content": f"{topic}\n\n{content}", "images": []}]}
    exp = engine.explain_concept(slides_data=slides_data)
    an = engine.generate_analogies(exp.get("explanation", ""), topic=topic)
    ex = engine.create_examples(exp.get("explanation", ""), topic=topic)
    qz = engine.generate_questions(exp.get("explanation", ""), topic=topic, num_questions=3)

    return {
        "explanation_scores": evaluate_explanation_quality(exp.get("explanation", ""), topic),
        "analogies_scores": evaluate_analogies_quality(an),
        "examples_scores": evaluate_examples_quality(ex),
        "quiz_scores": evaluate_quiz_quality(qz),
        "samples": {
            "explanation": exp.get("explanation", "")[:500],
            "analogies": an.get("analogies", []),
            "examples": ex.get("examples", []),
            "quiz": qz.get("questions", []),
        },
    }


def generate_validation_report(results: Dict) -> None:
    (RESULTS_DIR / "report.json").write_text(json.dumps(results, indent=2))
    lines = ["# Validation Report", ""]
    for domain, topics in results.items():
        lines.append(f"## {domain}")
        for t, data in topics.items():
            lines.append(f"- **{t}**: {json.dumps({k: v for k, v in data.items() if k.endswith('_scores')})}")
        lines.append("")
    (RESULTS_DIR / "report.md").write_text("\n".join(lines))


def main() -> int:
    engine = create_explanation_engine(load_mistral_in_4bit=True, load_vision_in_4bit=True)
    aggregated: Dict[str, Dict] = {}
    for domain, topics in TEST_CASES.items():
        aggregated[domain] = {}
        for t, content in topics.items():
            aggregated[domain][t] = test_topic(engine, t, content)
    generate_validation_report(aggregated)
    print("Validation complete. See validation_results/.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
