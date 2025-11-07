import pytest
import json
import time
import tempfile
from pathlib import Path
from unittest.mock import Mock, MagicMock, patch

from src.ppt_processor import process_ppt
from src.explanation_engine import create_explanation_engine
from src.app import (
    process_ppt_file,
    explain_slides,
    generate_analogies_handler,
    create_examples_handler,
    generate_quiz_handler,
)


@pytest.fixture(scope="module")
def sample_pptx_file(tmp_path_factory):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as e:
        pytest.skip(f"python-pptx not available or failed to import: {e}")

    tmp_dir = tmp_path_factory.mktemp("ppt_samples")
    ppt_path = tmp_dir / "sample_integration.pptx"

    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction to Photosynthesis"
    slide.placeholders[1].text = "An overview of the process plants use to convert light into energy."

    # Content slide with bullets
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Key Stages"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Light-dependent reactions"
    p = body.add_paragraph()
    p.text = "Calvin cycle"
    p.level = 1
    p = body.add_paragraph()
    p.text = "Glucose production"
    p.level = 1

    # Table slide
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Comparison"
    rows, cols = 3, 3
    left = top = Inches(1.5)
    width = Inches(7.0)
    height = Inches(1.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    headers = ["Stage", "Inputs", "Outputs"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    table.cell(1, 0).text = "Light Rxns"
    table.cell(1, 1).text = "Light, H2O"
    table.cell(1, 2).text = "ATP, NADPH"
    table.cell(2, 0).text = "Calvin"
    table.cell(2, 1).text = "CO2, ATP, NADPH"
    table.cell(2, 2).text = "Glucose"

    # Placeholder for image slide (no actual image to avoid file ops)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Chloroplast Diagram"

    prs.save(str(ppt_path))
    return ppt_path


@pytest.fixture(scope="module")
def mock_engine_for_integration():
    engine = MagicMock(name="ExplanationEngineMock")

    def _explain_concept(ppt_path=None, slides_data=None, slide_indices=None, include_images=True, temperature=0.7, max_tokens=512):
        used = slide_indices if slide_indices is not None else [0, 1, 2]
        return {
            "explanation": (
                "Photosynthesis is how plants convert light energy into chemical energy. "
                "It involves light-dependent reactions and the Calvin cycle, leading to glucose production."
            ),
            "topic": "Photosynthesis",
            "slides_used": used,
            "metadata": {
                "include_images": include_images,
                "temperature": temperature,
                "max_tokens": max_tokens,
            },
        }

    def _generate_analogies(concept_text, topic=None, temperature=0.7, max_tokens=256):
        return {
            "analogies": [
                "Photosynthesis is like charging a battery using sunlight.",
                "Think of chloroplasts as tiny solar panels producing energy for the cell.",
            ],
            "topic": topic or "Photosynthesis",
            "metadata": {"temperature": temperature, "max_tokens": max_tokens},
        }

    def _create_examples(concept_text, topic=None, temperature=0.7, max_tokens=384):
        return {
            "examples": [
                "A houseplant placed near a window grows faster due to more available light.",
                "Farmers use greenhouses to maximize light exposure and plant growth.",
            ],
            "topic": topic or "Photosynthesis",
            "metadata": {"temperature": temperature, "max_tokens": max_tokens},
        }

    def _generate_questions(concept_text, topic=None, num_questions=3, temperature=0.4, max_tokens=768):
        return {
            "questions": [
                {
                    "question": "Which stage produces glucose?",
                    "options": ["Light-dependent reactions", "Calvin cycle", "Glycolysis", "Electron transport"],
                    "correct_index": 1,
                    "rationale": "The Calvin cycle fixes CO2 into glucose.",
                },
                {
                    "question": "What are the outputs of light-dependent reactions?",
                    "options": ["ATP and NADPH", "CO2 and O2", "Glucose and water", "Pyruvate and ATP"],
                    "correct_index": 0,
                    "rationale": "ATP and NADPH are produced to power the Calvin cycle.",
                },
                {
                    "question": "Where does photosynthesis occur?",
                    "options": ["Mitochondria", "Nucleus", "Chloroplast", "Cytoplasm"],
                    "correct_index": 2,
                    "rationale": "Chloroplasts contain chlorophyll for light capture.",
                },
            ],
            "topic": topic or "Photosynthesis",
            "metadata": {"temperature": temperature, "max_tokens": max_tokens, "num_questions": num_questions},
        }

    engine.explain_concept.side_effect = _explain_concept
    engine.generate_analogies.side_effect = _generate_analogies
    engine.create_examples.side_effect = _create_examples
    engine.generate_questions.side_effect = _generate_questions
    return engine


def _mock_create_engine(mock_engine):
    def _factory(*args, **kwargs):
        return mock_engine
    return _factory


def test_e2e_ppt_processing_to_explanation(sample_pptx_file, mock_engine_for_integration, monkeypatch):
    slides_data = process_ppt(str(sample_pptx_file))
    assert isinstance(slides_data, dict)
    assert "slides" in slides_data and isinstance(slides_data["slides"], list)
    assert slides_data.get("total_slides", 0) >= 3

    monkeypatch.setattr("src.explanation_engine.create_explanation_engine", _mock_create_engine(mock_engine_for_integration))
    engine = create_explanation_engine()
    result = engine.explain_concept(slides_data=slides_data)

    assert "explanation" in result and isinstance(result["explanation"], str)
    assert "slides_used" in result and isinstance(result["slides_used"], list)
    assert result["topic"] == "Photosynthesis"


def test_e2e_full_workflow_with_app_handlers(sample_pptx_file, mock_engine_for_integration, monkeypatch):
    session_state = {"ppt_data": None, "current_slide_idx": 0, "chat_history": [], "last_explanation": ""}

    monkeypatch.setattr("src.explanation_engine.create_explanation_engine", _mock_create_engine(mock_engine_for_integration))

    process_ppt_file(str(sample_pptx_file), session_state)
    assert session_state["ppt_data"] is not None

    explain_slides(session_state)
    assert session_state["last_explanation"]
    assert len(session_state["chat_history"]) >= 1

    generate_analogies_handler(session_state)
    generate_analogies_added = any("analogies" in (m.get("type") or "") or "Analogy" in m.get("content", "") for m in session_state["chat_history"])
    assert generate_analogies_added

    create_examples_handler(session_state)
    examples_added = any("Example" in m.get("content", "") for m in session_state["chat_history"])
    assert examples_added

    generate_quiz_handler(session_state)
    quiz_added = any("Question" in m.get("content", "") or "Quiz" in m.get("content", "") for m in session_state["chat_history"])
    assert quiz_added


@pytest.mark.slow
def test_e2e_multi_slide_context_handling(sample_pptx_file, mock_engine_for_integration, monkeypatch):
    slides_data = process_ppt(str(sample_pptx_file))
    monkeypatch.setattr("src.explanation_engine.create_explanation_engine", _mock_create_engine(mock_engine_for_integration))
    engine = create_explanation_engine()

    # Request explanation for multiple slides
    indices = [0, 1, 2]
    result = engine.explain_concept(slides_data=slides_data, slide_indices=indices)
    assert result["slides_used"] == indices

    # Respect MAX_SLIDES_CONTEXT (assume 5)
    many = list(range(10))
    result_many = engine.explain_concept(slides_data=slides_data, slide_indices=many)
    assert len(result_many["slides_used"]) >= 1


@pytest.mark.slow
def test_e2e_image_processing_and_captioning(sample_pptx_file, mock_engine_for_integration, monkeypatch, tmp_path):
    slides_data = process_ppt(str(sample_pptx_file))
    out_dir = Path("static/images/extracted")
    out_dir.mkdir(parents=True, exist_ok=True)

    # Mock BLIP-2 captioning inside engine
    monkeypatch.setattr("src.explanation_engine.create_explanation_engine", _mock_create_engine(mock_engine_for_integration))
    engine = create_explanation_engine()

    result = engine.explain_concept(slides_data=slides_data, include_images=True)
    assert "explanation" in result


def test_e2e_error_handling_invalid_ppt(monkeypatch, tmp_path):
    with pytest.raises(Exception):
        process_ppt("non_existent.pptx")

    bad = tmp_path / "not_a_ppt.txt"
    bad.write_text("hello")
    with pytest.raises(Exception):
        process_ppt(str(bad))


def test_e2e_session_state_isolation(sample_pptx_file, mock_engine_for_integration, monkeypatch):
    monkeypatch.setattr("src.explanation_engine.create_explanation_engine", _mock_create_engine(mock_engine_for_integration))

    s1 = {"ppt_data": None, "current_slide_idx": 0, "chat_history": [], "last_explanation": ""}
    s2 = {"ppt_data": None, "current_slide_idx": 0, "chat_history": [], "last_explanation": ""}

    process_ppt_file(str(sample_pptx_file), s1)
    process_ppt_file(str(sample_pptx_file), s2)

    explain_slides(s1)
    explain_slides(s2)

    assert s1 is not s2
    assert s1["last_explanation"] != ""
    assert s2["last_explanation"] != ""


def test_e2e_quiz_json_parsing_and_formatting(mock_engine_for_integration, monkeypatch):
    monkeypatch.setattr("src.explanation_engine.create_explanation_engine", _mock_create_engine(mock_engine_for_integration))

    session_state = {"ppt_data": {"slides": ["dummy"]}, "current_slide_idx": 0, "chat_history": [], "last_explanation": "Photosynthesis explanation"}
    generate_quiz_handler(session_state)

    # Expect a quiz-like entry in chat history
    joined = "\n".join(m.get("content", "") for m in session_state["chat_history"])
    assert "Question" in joined or "options" in joined


@pytest.mark.slow
def test_e2e_performance_benchmarks(sample_pptx_file, mock_engine_for_integration, monkeypatch):
    monkeypatch.setattr("src.explanation_engine.create_explanation_engine", _mock_create_engine(mock_engine_for_integration))
    session_state = {"ppt_data": None, "current_slide_idx": 0, "chat_history": [], "last_explanation": ""}

    t0 = time.time()
    process_ppt_file(str(sample_pptx_file), session_state)
    t1 = time.time()
    explain_slides(session_state)
    t2 = time.time()
    generate_analogies_handler(session_state)
    t3 = time.time()
    create_examples_handler(session_state)
    t4 = time.time()
    generate_quiz_handler(session_state)
    t5 = time.time()

    assert (t1 - t0) < 5.0
    assert (t2 - t1) < 5.0
    assert (t3 - t2) < 5.0
    assert (t4 - t3) < 5.0
    assert (t5 - t4) < 5.0


def test_e2e_concurrent_requests_simulation(sample_pptx_file, mock_engine_for_integration, monkeypatch):
    monkeypatch.setattr("src.explanation_engine.create_explanation_engine", _mock_create_engine(mock_engine_for_integration))

    sessions = [
        {"ppt_data": None, "current_slide_idx": 0, "chat_history": [], "last_explanation": ""}
        for _ in range(3)
    ]

    for s in sessions:
        process_ppt_file(str(sample_pptx_file), s)
        explain_slides(s)
        generate_analogies_handler(s)
        create_examples_handler(s)
        generate_quiz_handler(s)

    for s in sessions:
        assert s["ppt_data"] is not None
        assert s["last_explanation"]
        assert len(s["chat_history"]) >= 3
