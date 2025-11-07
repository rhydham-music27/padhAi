import json
from pathlib import Path

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.ppt_processor import _iter_shapes, extract_images, extract_text, process_ppt


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    prs = Presentation()

    # Slide 1: Title slide
    title_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Test Title"
    slide1.placeholders[1].text = "Test Subtitle"

    # Slide 2: Content with bullets and a 2x2 table
    layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(layout)
    body = slide2.shapes.placeholders[1].text_frame
    body.clear()
    p1 = body.paragraphs[0]
    p1.text = "Bullet 1"
    p1.level = 0
    p2 = body.add_paragraph()
    p2.text = "Bullet 2"
    p2.level = 1

    # Add table 2x2
    rows, cols = 2, 2
    left, top, width, height = Inches(1), Inches(3), Inches(6), Inches(1)
    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "R1C1"
    table.cell(0, 1).text = "R1C2"
    table.cell(1, 0).text = "R2C1"
    table.cell(1, 1).text = "R2C2"

    # Slide 3: Picture
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])

    # Create a simple image with Pillow
    img_path = tmp_path / "test_img.png"
    img = PILImage.new("RGB", (120, 60), color=(10, 200, 100))
    img.save(img_path)

    slide3.shapes.add_picture(str(img_path), Inches(1), Inches(1), width=Inches(2))

    file_path = tmp_path / "test_presentation.pptx"
    prs.save(file_path)
    return file_path


@pytest.fixture
def output_dir(tmp_path: Path) -> Path:
    out = tmp_path / "extracted_images"
    return out


def test_extract_text_basic(sample_pptx: Path):
    data = extract_text(str(sample_pptx))
    assert isinstance(data, list)
    assert len(data) == 3
    # Slide numbers 1..3
    assert [d["slide_number"] for d in data] == [1, 2, 3]
    # First slide should contain title text
    assert "Test Title" in data[0]["text_content"]
    # Required keys
    for slide in data:
        for key in ["slide_number", "text_content", "tables", "notes", "shape_details"]:
            assert key in slide


def test_extract_text_tables(sample_pptx: Path):
    data = extract_text(str(sample_pptx))
    slide2 = data[1]
    assert slide2["tables"], "Expected tables on slide 2"
    tbl = slide2["tables"][0]
    assert len(tbl) == 2 and len(tbl[0]) == 2
    assert tbl[0][0] == "R1C1"
    assert tbl[1][1] == "R2C2"


def test_extract_text_empty_slide(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    empty_path = tmp_path / "empty.pptx"
    prs.save(empty_path)

    data = extract_text(str(empty_path))
    assert len(data) == 1
    assert data[0]["text_content"] == ""
    assert data[0]["tables"] == []


def test_extract_text_file_not_found():
    with pytest.raises(FileNotFoundError):
        extract_text("nonexistent.pptx")


def test_extract_images_basic(sample_pptx: Path, output_dir: Path):
    imgs = extract_images(str(sample_pptx), str(output_dir))
    assert isinstance(imgs, list)
    assert imgs, "Expected at least one image extracted"
    for meta in imgs:
        for key in ["slide_number", "image_path", "image_format", "sha1", "shape_name"]:
            assert key in meta
        assert Path(meta["image_path"]).exists()
        # extension matches format
        assert Path(meta["image_path"]).suffix.replace(".", "") == meta["image_format"].replace(".", "")


def test_extract_images_deduplication(tmp_path: Path, output_dir: Path):
    # Create PPT with same image on multiple slides
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])

    img_path = tmp_path / "dup.png"
    PILImage.new("RGB", (64, 64), color=(255, 0, 0)).save(img_path)

    slide1.shapes.add_picture(str(img_path), Inches(1), Inches(1), width=Inches(1))
    slide2.shapes.add_picture(str(img_path), Inches(2), Inches(2), width=Inches(1))

    file_path = tmp_path / "dup.pptx"
    prs.save(file_path)

    metas = extract_images(str(file_path), str(output_dir))
    assert len(metas) >= 2
    sha1s = [m["sha1"] for m in metas]
    assert len(set(sha1s)) == 1, "Images should deduplicate by SHA1"
    # Only one physical file saved
    saved_files = list(Path(output_dir).glob("*"))
    assert len(saved_files) == 1


def test_extract_images_no_images(tmp_path: Path, output_dir: Path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "noimg.pptx"
    prs.save(path)

    metas = extract_images(str(path), str(output_dir))
    assert metas == []
    assert not Path(output_dir).exists() or not any(Path(output_dir).glob("*"))


def test_process_ppt_full(sample_pptx: Path, output_dir: Path):
    result = process_ppt(str(sample_pptx), extract_images_flag=True, output_dir=str(output_dir))
    assert set(result.keys()) == {"file_path", "total_slides", "slides", "images", "processing_errors"}
    assert result["total_slides"] == 3
    assert len(result["slides"]) == 3
    # Each slide should have images key
    for s in result["slides"]:
        assert "images" in s
    assert len(result["images"]) >= 1
    assert result["processing_errors"] == []


def test_process_ppt_text_only(sample_pptx: Path, output_dir: Path):
    result = process_ppt(str(sample_pptx), extract_images_flag=False, output_dir=str(output_dir))
    assert result["images"] == []
    assert not Path(output_dir).exists() or not any(Path(output_dir).glob("*"))
    assert len(result["slides"]) == 3


def test_process_ppt_invalid_file():
    with pytest.raises(FileNotFoundError):
        process_ppt("invalid.pptx")


def test_iter_shapes_grouped():
    # Simulate grouped shapes using simple mocks to validate recursion logic
    class MockShape:
        def __init__(self, text: str = "", name: str = "TextBox"):
            self.name = name
            self.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
            self.has_text_frame = True

            class TF:
                def __init__(self, t: str):
                    self.paragraphs = []
                    class P:
                        def __init__(self, t: str):
                            self.text = t
                            self.runs = []
                    self.paragraphs.append(P(t))
            self.text_frame = TF(text)

    class MockGroup:
        def __init__(self, shapes):
            self.shape_type = MSO_SHAPE_TYPE.GROUP
            self.shapes = shapes

    inner = [MockShape("Inside Group")] 
    group = MockGroup(inner)
    top_level = [group]

    gathered = list(_iter_shapes(top_level))
    assert any(getattr(s, "has_text_frame", False) for s in gathered)
