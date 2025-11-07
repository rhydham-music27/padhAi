"""
PPT processing module for extracting text and images from PowerPoint (.pptx) files.

Functions:
- _iter_shapes: recursively yield shapes, including inside groups
- extract_text: extract text, tables, and notes per slide
- extract_images: extract images with deduplication and metadata
- process_ppt: unified entry that merges text and images by slide
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, Generator, List, Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


def _iter_shapes(shapes) -> Generator[Any, None, None]:
    """Recursively yield shapes from a slide, descending into groups.

    Args:
        shapes: A shape collection, e.g., slide.shapes

    Yields:
        Individual shape objects, including those nested within grouped shapes.
    """
    for shp in shapes:
        try:
            if getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                # Descend into grouped shapes
                yield from _iter_shapes(shp.shapes)
            else:
                yield shp
        except Exception as exc:  # defensive: corrupted shapes shouldn't stop iteration
            logger.warning("Skipping shape due to error: %s", exc)
            continue


def extract_text(pptx_path: str) -> List[Dict[str, Any]]:
    """Extract text, tables, and notes from a PPTX file per slide.

    Args:
        pptx_path: Path to the .pptx file

    Returns:
        A list of dictionaries, each with keys:
            - slide_number (int)
            - text_content (str)
            - tables (List[List[List[str]]])  # list of 2D tables
            - notes (Optional[str])
            - shape_details (List[Dict[str, Any]])
    """
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError("Only .pptx files are supported")

    results: List[Dict[str, Any]] = []
    try:
        prs = Presentation(str(path))
        logger.info("Loaded presentation: %s", path)
    except Exception as exc:
        logger.error("Failed to load presentation %s: %s", path, exc)
        raise

    for idx, slide in enumerate(prs.slides, start=1):
        text_items: List[str] = []
        tables: List[List[List[str]]] = []
        shape_details: List[Dict[str, Any]] = []
        notes_text: Optional[str] = None

        try:
            for shp in _iter_shapes(slide.shapes):
                # Text frames
                if hasattr(shp, "has_text_frame") and getattr(shp, "has_text_frame", False):
                    try:
                        parts: List[str] = []
                        tf = shp.text_frame
                        if tf is not None:
                            for p in tf.paragraphs:
                                runs_text = "".join(run.text for run in p.runs) if p.runs else p.text
                                parts.append(runs_text)
                        text_value = "\n".join(t for t in parts if t is not None)
                        if text_value:
                            text_items.append(text_value)
                        shape_details.append({
                            "shape_name": getattr(shp, "name", ""),
                            "text": text_value,
                            "shape_type": str(getattr(shp, "shape_type", "")),
                        })
                    except Exception as exc:
                        logger.warning("Failed extracting text from shape on slide %d: %s", idx, exc)

                # Tables
                if hasattr(shp, "has_table") and getattr(shp, "has_table", False):
                    try:
                        tbl = shp.table
                        table_matrix: List[List[str]] = []
                        for row in tbl.rows:
                            row_vals: List[str] = []
                            for cell in row.cells:
                                # Prefer cell.text which concatenates paragraphs
                                row_vals.append(cell.text or "")
                            table_matrix.append(row_vals)
                        tables.append(table_matrix)
                    except Exception as exc:
                        logger.warning("Failed extracting table from slide %d: %s", idx, exc)

            # Speaker notes
            try:
                if getattr(slide, "has_notes_slide", False) and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text or None
            except Exception as exc:
                logger.warning("Failed extracting notes from slide %d: %s", idx, exc)

            slide_text = "\n".join(text_items).strip()
            results.append({
                "slide_number": idx,
                "text_content": slide_text,
                "tables": tables,
                "notes": notes_text,
                "shape_details": shape_details,
            })
        except Exception as exc:
            logger.error("Error processing slide %d: %s", idx, exc)
            results.append({
                "slide_number": idx,
                "text_content": "",
                "tables": [],
                "notes": None,
                "shape_details": [],
            })

    return results


def extract_images(pptx_path: str, output_dir: str = "static/images/extracted") -> List[Dict[str, Any]]:
    """Extract embedded images from PPTX, save to disk, and return metadata.

    Args:
        pptx_path: Path to .pptx file
        output_dir: Directory to save extracted images

    Returns:
        List of dictionaries with keys:
            - slide_number (int)
            - image_path (str)
            - image_format (str)
            - sha1 (str)
            - shape_name (str)
            - dimensions (Optional[Dict[str, int]])
    """
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError("Only .pptx files are supported")

    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    try:
        prs = Presentation(str(path))
        logger.info("Loaded presentation for images: %s", path)
    except Exception as exc:
        logger.error("Failed to load presentation %s for images: %s", path, exc)
        raise

    seen_hashes: Dict[str, Path] = {}
    image_metadata: List[Dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        for shp in _iter_shapes(slide.shapes):
            try:
                if getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    img = shp.image
                    sha1 = img.sha1
                    ext = img.ext  # e.g., 'png', 'jpeg'
                    file_ext = f".{ext}" if not str(ext).startswith(".") else str(ext)
                    if sha1 in seen_hashes:
                        file_path = seen_hashes[sha1]
                    else:
                        filename = f"slide{idx:03d}_{sha1}{file_ext}"
                        file_path = out_dir / filename
                        try:
                            file_path.write_bytes(img.blob)
                            seen_hashes[sha1] = file_path
                        except Exception as exc:
                            logger.warning("Failed to save image on slide %d: %s", idx, exc)
                            continue

                    dims: Optional[Dict[str, int]] = None
                    try:
                        with Image.open(file_path) as pil_img:
                            w, h = pil_img.size
                            dims = {"width": int(w), "height": int(h)}
                    except Exception:
                        # Unsupported formats (e.g., EMF/WMF) may not be readable by Pillow
                        dims = None

                    image_metadata.append({
                        "slide_number": idx,
                        "image_path": str(file_path),
                        "image_format": ext,
                        "sha1": sha1,
                        "shape_name": getattr(shp, "name", ""),
                        "dimensions": dims,
                    })
            except Exception as exc:
                logger.warning("Skipping picture on slide %d due to error: %s", idx, exc)
                continue

    return image_metadata


def process_ppt(
    pptx_path: str,
    extract_images_flag: bool = True,
    output_dir: str = "static/images/extracted",
) -> Dict[str, Any]:
    """Unified processing of a PPTX file to extract text and images.

    Args:
        pptx_path: Path to the .pptx file
        extract_images_flag: Whether to extract images in addition to text
        output_dir: Directory to store extracted images

    Returns:
        Dictionary with keys:
            - file_path (str)
            - total_slides (int)
            - slides (List[Dict])
            - images (List[Dict])
            - processing_errors (List[str])
    """
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError("Only .pptx files are supported")

    errors: List[str] = []

    # Text extraction
    try:
        text_data = extract_text(str(path))
    except Exception as exc:
        logger.error("Text extraction failed: %s", exc)
        errors.append(f"text_extraction_failed: {exc}")
        text_data = []

    # Image extraction
    images: List[Dict[str, Any]] = []
    if extract_images_flag:
        try:
            images = extract_images(str(path), output_dir=output_dir)
        except Exception as exc:
            logger.error("Image extraction failed: %s", exc)
            errors.append(f"image_extraction_failed: {exc}")
            images = []

    # Merge per slide
    images_by_slide: Dict[int, List[Dict[str, Any]]] = {}
    for meta in images:
        images_by_slide.setdefault(int(meta.get("slide_number", 0)), []).append(meta)

    slides_out: List[Dict[str, Any]] = []
    for slide_info in text_data:
        sn = int(slide_info.get("slide_number", 0))
        slide_images = images_by_slide.get(sn, [])
        merged = dict(slide_info)
        merged["images"] = slide_images
        slides_out.append(merged)

    result: Dict[str, Any] = {
        "file_path": str(path),
        "total_slides": len(text_data),
        "slides": slides_out,
        "images": images,
        "processing_errors": errors,
    }

    return result
